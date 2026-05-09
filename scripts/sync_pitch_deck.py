#!/usr/bin/env python3
"""
sync_pitch_deck.py — NouchiX / SouHimBou AI Pitch Deck Auto-Updater

Patches SouHimBouAI_PitchDeck_NouchiX.pptx in place with:
  1. Live KPI figures pulled from the Supabase biz_matrix_snapshots table
  2. Static narrative improvements:
       - Brave Search ZDR integration added to capabilities, tech stack,
         competitive advantage table, and traction slide
       - Offer ladder reconciled ($19/mo Beta -> $60k-$120k Pilot -> $250k+ Enterprise)
       - Deployment architecture tiering (VPS / Fly.io / Vercel / Supabase / Brave ZDR)

Usage:
    pip install python-pptx requests
    export SUPABASE_URL=...
    export SUPABASE_SERVICE_ROLE_KEY=...
    python scripts/sync_pitch_deck.py [--deck path/to/deck.pptx] [--dry-run]

In CI (GitHub Actions) the script reads env vars automatically.
"""

import argparse
import os
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

try:
    import requests
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    sys.exit(
        "ERROR: Missing dependencies. Run: pip install python-pptx requests"
    )


# --- KPI Fetch ----------------------------------------------------------------

def fetch_latest_kpis() -> dict[str, Any]:
    """Pull the most recent biz_matrix_snapshots row from Supabase."""
    url = os.environ.get("SUPABASE_URL", "")
    key = os.environ.get("SUPABASE_SERVICE_ROLE_KEY", "")
    if not url or not key:
        print("[WARN] SUPABASE_URL / SUPABASE_SERVICE_ROLE_KEY not set -- using defaults")
        return {}

    resp = requests.get(
        f"{url}/rest/v1/biz_matrix_snapshots",
        headers={
            "apikey": key,
            "Authorization": f"Bearer {key}",
            "Range": "0-0",
        },
        params={"order": "snapshot_at.desc", "limit": "1"},
        timeout=15,
    )
    if not resp.ok:
        print(f"[WARN] KPI fetch failed ({resp.status_code}) -- using defaults")
        return {}

    rows = resp.json()
    return rows[0] if rows else {}


# --- Text Replacement Engine --------------------------------------------------

def replace_text_in_shape(shape: Any, old: str, new: str) -> bool:
    changed = False
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                changed = True
    return changed


def apply_replacements(prs: Presentation, replacements: list[tuple[str, str]]) -> int:
    total = 0
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            for old, new in replacements:
                if replace_text_in_shape(shape, old, new):
                    total += 1
                    print(f"  [Slide {slide_idx}] '{old[:60]}' -> '{new[:60]}'")
    return total


# --- Static Narrative Patches -------------------------------------------------

STATIC_NARRATIVE_PATCHES: list[tuple[str, str]] = [
    # Slide 5: capabilities
    (
        "Real-time continuous monitoring",
        "Real-time continuous monitoring | Real-time cited DISA intelligence (Brave Search ZDR)"
    ),
    # Slide 7: differentiation checklist
    (
        "Defense Alignment (CMMC/STIG)",
        "Defense Alignment (CMMC/STIG)\nReal-Time DISA Grounding (Brave ZDR)\nPrivacy-Preserving AI (No Query Retention)"
    ),
    # Slide 8: tech stack
    (
        "Tehama secure VDI",
        "Tehama secure VDI | Brave Search API (Zero Data Retention)"
    ),
    # Slide 9: offer ladder
    (
        "$80K–$200K per facility/year",
        "Beta $19/mo -> Pilot $60K-$120K -> Enterprise $250K+"
    ),
    # Slide 10: competitive advantage — add ZDR rows note
    (
        "Veteran-Led Expertise",
        "Veteran-Led Expertise | Real-Time DISA Intel (Brave ZDR) | Zero Data Retention AI"
    ),
    # Slide 12: roadmap
    (
        "Pilot Program LOIs",
        "Pilot Program LOIs | Brave ZDR Real-Time DISA Grounding (Sprint 36)"
    ),
    # Slide 16: data handling statement
    (
        "ai-nativevc@souhimbou.ai",
        "Data Handling: Brave ZDR -- no query data retained outside your org. | ai-nativevc@souhimbou.ai"
    ),
]


def build_dynamic_patches(kpis: dict[str, Any]) -> list[tuple[str, str]]:
    patches: list[tuple[str, str]] = []

    dod = kpis.get("dod_contractors_secured")
    if dod is not None:
        patches.append(("47 defense contractors secured", f"{dod} defense contractors secured"))

    scans = kpis.get("compliance_scans_7d")
    if scans is not None:
        patches.append(("500+ defense systems secured", f"{scans}+ compliance scans (7d)"))

    brave_hits = kpis.get("brave_grounding_hits_7d")
    if brave_hits is not None:
        patches.append((
            "< 15 min threat response time",
            f"< 15 min threat response  |  {brave_hits} Brave-grounded answers (7d)"
        ))

    snap = kpis.get("snapshot_at", datetime.now(timezone.utc).isoformat())
    date_str = snap[:10]
    patches.append((
        "STIG-First Compliance Autopilot MVP",
        f"STIG-First Compliance Autopilot MVP | Last synced: {date_str}"
    ))

    return patches


# --- Main --------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Sync NouchiX Pitch Deck")
    parser.add_argument(
        "--deck",
        default=str(
            Path(__file__).parent.parent / "assets" / "SouHimBouAI_PitchDeck_NouchiX.pptx"
        ),
        help="Path to the source .pptx file",
    )
    parser.add_argument("--out", default=None, help="Output path")
    parser.add_argument("--dry-run", action="store_true", help="Preview without writing")
    args = parser.parse_args()

    deck_path = Path(args.deck)
    if not deck_path.exists():
        candidates = list(Path(".").rglob("SouHimBouAI_PitchDeck_NouchiX.pptx"))
        if candidates:
            deck_path = candidates[0]
            print(f"[INFO] Found deck at {deck_path}")
        else:
            sys.exit(f"ERROR: Deck not found at {args.deck}. Pass --deck <path>.")

    out_path = Path(args.out) if args.out else deck_path
    print(f"\n{'[DRY-RUN] ' if args.dry_run else ''}Syncing: {deck_path}")
    print("-" * 60)

    prs = Presentation(str(deck_path))

    print("\n[1/3] Fetching live KPIs from Supabase...")
    kpis = fetch_latest_kpis()
    if kpis:
        print(f"      Snapshot: {kpis.get('snapshot_at', 'N/A')}")
        print(f"      DoD contractors: {kpis.get('dod_contractors_secured', 'N/A')}")
        print(f"      Brave hits (7d): {kpis.get('brave_grounding_hits_7d', 'N/A')}")
    else:
        print("      No live data -- static narrative patches only")

    all_patches = STATIC_NARRATIVE_PATCHES + build_dynamic_patches(kpis)
    print(f"\n[2/3] Applying {len(all_patches)} patch(es)...")
    total_changes = apply_replacements(prs, all_patches)
    print(f"      {total_changes} text run(s) updated")

    if args.dry_run:
        print("\n[3/3] DRY-RUN -- no file written")
    else:
        prs.save(str(out_path))
        print(f"\n[3/3] Saved -> {out_path}")

    print("\nDone.")


if __name__ == "__main__":
    main()
